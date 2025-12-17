"""
Generate a PowerPoint presentation with each graph as its own slide
from the augmented Survicate cancelled subscriptions data.
"""

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import sys
from io import BytesIO

def generate_churn_slides(csv_path='data/survicate_cancelled_subscriptions_augmented.csv', 
                          output_path='churn_trends_slides.pptx'):
    """
    Generate a PowerPoint presentation with each graph as its own slide
    
    Args:
        csv_path: Path to the augmented CSV file
        output_path: Path where the PowerPoint will be saved
    """
    # Check if CSV exists
    if not os.path.exists(csv_path):
        print(f"Error: CSV file not found at {csv_path}")
        print(f"Current working directory: {os.getcwd()}")
        return None
    
    print(f"Reading data from: {csv_path}")
    df = pd.read_csv(csv_path)
    
    print(f"Total rows: {len(df)}")
    
    # Filter out rows with missing data
    df = df[df['augmented_churn_reason'].notna() & df['year_month'].notna()]
    df = df[df['year_month'] != '2024-11']  # Exclude November 2024
    
    print(f"Rows with valid churn reason and year_month: {len(df)}")
    
    if len(df) == 0:
        print("Error: No valid data found after filtering")
        return None
    
    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color palette (same as generate_churn_report.py)
    colors = [
        '#4285F4', '#EA4335', '#FBBC04', '#34A853', '#FF6D01', '#9334E6', '#00ACC1',
        '#64B5F6', '#F28B82', '#FFF176', '#81C784', '#FFB74D', '#BA68C8', '#4DB6AC',
        '#BBDEFB', '#FAD2CF', '#FFF9C4', '#C8E6C9', '#FFE0B2', '#E1BEE7', '#B2DFDB',
        '#E8F0FE', '#FCE8E6',
    ]
    
    # SLIDE 1: Main Churn Trends Chart
    print("Generating main churn trends chart...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Churn Reasons Over Time"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Generate main chart
    chart_img = create_main_chart(df, colors)
    if chart_img:
        slide.shapes.add_picture(chart_img, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    
    # SLIDE 2-N: Question Trend Charts
    # Updated to reflect new question order after Q9/Q10 were added
    question_mapping = {
        'Q2': 'Q#2: Where does the location pin not match your dog\'s location?',
        'Q3': 'Q#3: Was the pet location pin grayed out when the location was inaccurate?',
        'Q4': 'Q#4: Is the collar not sending feedback or is your dog not responding to the feedback sent?',
        'Q5': 'Q#5: Did you screw in the contact tips required for static feedback to work properly?',
        'Q6': 'Q#6: What battery life, charging or power issues did you encounter?',
        'Q7': 'Q#7: Which containment solution did you purchase?',
        'Q8': 'Q#8: Which other GPS or wireless fence did you purchase?',
        'Q11': 'Q#11: Did you engage with the Learn training curriculum?',  # Was Q9, now Q11
        'Q12': 'Q#12: What was the main reason you didn\'t complete the Learn curriculum?',  # Was Q10, now Q12
        'Q13': 'Q#13: Did you contact our Customer Service team via Dog Park?',  # Was Q11, now Q13
        'Q14': 'Q#14: Would a free session with a trainer to help your dog use the collar effectively have helped you continue to use it?'  # Was Q12, now Q14
    }
    
    for question_id, question_text in question_mapping.items():
        print(f"Generating chart for {question_id}...")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = question_text
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.LEFT
        
        # Generate question chart
        chart_img = create_question_chart(df, question_id, question_text)
        if chart_img:
            slide.shapes.add_picture(chart_img, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    
    # Save presentation
    print(f"Generating PowerPoint: {output_path}")
    prs.save(output_path)
    print(f"[OK] Presentation generated successfully: {output_path}")
    print(f"  - Total slides: {len(prs.slides)}")
    return output_path


def create_main_chart(df, colors):
    """Create the main churn trends chart and return as image bytes"""
    # Group by year_month and augmented_churn_reason
    grouped = df.groupby(['year_month', 'augmented_churn_reason']).size().reset_index(name='count')
    monthly_totals = df.groupby('year_month').size()
    grouped['percentage'] = grouped.apply(
        lambda row: (row['count'] / monthly_totals[row['year_month']]) * 100, 
        axis=1
    )
    
    # Pivot data
    pivot_data = grouped.pivot(index='augmented_churn_reason', columns='year_month', values='percentage').fillna(0)
    months = sorted(pivot_data.columns)
    pivot_data = pivot_data[months]
    
    # Sort reasons by total percentage
    reason_totals = pivot_data.sum(axis=1).sort_values(ascending=False)
    top_n = 11
    top_reasons = reason_totals.head(top_n).index.tolist()
    other_reasons = reason_totals.tail(len(reason_totals) - top_n).index.tolist() if len(reason_totals) > top_n else []
    
    # Aggregate "Other" if needed
    if len(other_reasons) > 0:
        other_data = pivot_data.loc[other_reasons].sum(axis=0)
        pivot_data = pivot_data.loc[top_reasons]
        pivot_data.loc['Other'] = other_data
    else:
        pivot_data = pivot_data.loc[top_reasons]
    
    # Re-sort
    reason_totals = pivot_data.sum(axis=1).sort_values(ascending=False)
    if 'Other' in reason_totals.index:
        other_total = reason_totals['Other']
        reason_totals = reason_totals.drop('Other')
        reason_totals['Other'] = other_total
    pivot_data = pivot_data.loc[reason_totals.index]
    
    # Create figure
    fig, ax = plt.subplots(figsize=(16, 10))
    
    # Create stacked bar chart
    bottom = None
    for i, (reason, row) in enumerate(pivot_data.iterrows()):
        ax.bar(months, row.values, bottom=bottom, 
              label=reason, color=colors[i % len(colors)], alpha=0.9, 
              edgecolor='white', linewidth=1.0)
        if bottom is None:
            bottom = row.values
        else:
            bottom = bottom + row.values
    
    # Formatting
    ax.set_xlabel('Year Month', fontsize=13, fontweight='600', color='#333333')
    ax.set_ylabel('Percentage of Churn (%)', fontsize=13, fontweight='600', color='#333333')
    ax.set_title('Churn Reasons Over Time', fontsize=17, fontweight='600', pad=20, color='#1a1a1a')
    ax.set_ylim(0, 100)
    ax.set_yticks([0, 20, 40, 60, 80, 100])
    ax.set_yticklabels(['0%', '20%', '40%', '60%', '80%', '100%'])
    ax.grid(axis='y', alpha=0.2, linestyle='--', color='#cccccc')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#e0e0e0')
    ax.spines['bottom'].set_color('#e0e0e0')
    
    plt.xticks(rotation=45, ha='right', fontsize=10, color='#666666')
    plt.yticks(fontsize=10, color='#666666')
    
    # Legend
    legend = ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=8.5, 
                       framealpha=0.95, frameon=True, edgecolor='#e0e0e0', 
                       facecolor='white', title='Churn Reasons', 
                       title_fontsize=9.5)
    legend.get_frame().set_linewidth(0.5)
    
    plt.tight_layout()
    
    # Convert to image bytes
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    
    return img_buffer


def create_question_chart(df, question_id, question_text):
    """Create a question trend chart and return as image bytes"""
    # Map question codes to search patterns (same as API route)
    # Updated to reflect new question order - use centralized config if available
    try:
        from backend.utils.survicate_question_config import get_question_by_legacy_number
        question_config = get_question_by_legacy_number(question_id)
        if question_config:
            pattern = question_config['patterns']
            exclude_patterns = question_config.get('exclude_patterns', [])
        else:
            # Fallback to hardcoded patterns
            pattern = []
            exclude_patterns = []
    except ImportError:
        # Fallback if config not available
        question_patterns = {
            'Q2': ['location pin', 'not match', 'dog'],
            'Q3': ['pet location pin', 'grayed out', 'inaccurate'],
            'Q4': ['collar', 'not sending feedback', 'dog not responding'],
            'Q5': ['screw in', 'contact tips', 'static feedback'],
            'Q6': ['battery life', 'charging', 'power issues'],
            'Q7': ['containment solution', 'purchase'],
            'Q8': ['other GPS', 'wireless fence', 'purchase'],
            'Q11': ['engage', 'Learn', 'training curriculum'],  # Was Q9, now Q11
            'Q12': ['main reason', 'didn\'t complete', 'Learn curriculum'],  # Was Q10, now Q12
            'Q13': ['contact', 'Customer Service', 'Dog Park'],  # Was Q11, now Q13
            'Q14': ['free session', 'trainer', 'collar effectively']  # Was Q12, now Q14
        }
        pattern = question_patterns.get(question_id, [])
        exclude_patterns = []
    
    # Find matching column (same logic as API route)
    pattern = question_patterns.get(question_id, [])
    column_name = None
    
    # First try pattern matching
    for col in df.columns:
        col_lower = str(col).lower()
        
        # Skip columns that match exclude patterns
        if exclude_patterns and any(exclude_word.lower() in col_lower for exclude_word in exclude_patterns):
            continue
        
        # Match by pattern (text content) - don't require Q# number match
        # This makes it resilient to question renumbering
        if pattern and all(word.lower() in col_lower for word in pattern):
            # Prefer columns with (Answer) or (Comment) suffix
            if '(Answer)' in col or '(Comment)' in col:
                column_name = col
                break
    
    # If pattern matching didn't work, try exact match
    if not column_name:
        for col in df.columns:
            col_str = str(col)
            if col_str.startswith(f'Q#{question_id[1:]}:') or col_str.startswith(f'Q{question_id[1:]}:'):
                if question_id in ['Q4', 'Q6', 'Q9']:
                    if '(Answer)' in col_str:
                        column_name = col
                        break
                else:
                    column_name = col
                    break
    
    if not column_name:
        print(f"Warning: Could not find column for {question_id}")
        return None
    
    # Filter to only rows with responses to this question
    df_filtered = df[df[column_name].notna() & (df[column_name].astype(str).str.strip() != '')].copy()
    
    if len(df_filtered) == 0:
        print(f"Warning: No responses found for {question_id}")
        return None
    
    # Get unique answers
    df_filtered['answer'] = df_filtered[column_name].astype(str).str.strip()
    unique_answers = df_filtered['answer'].unique().tolist()
    
    # Group by month and answer
    grouped = df_filtered.groupby(['year_month', 'answer']).size().reset_index(name='count')
    monthly_totals = df_filtered.groupby('year_month').size()
    grouped['percentage'] = grouped.apply(
        lambda row: (row['count'] / monthly_totals[row['year_month']]) * 100,
        axis=1
    )
    
    # Calculate total counts for each answer (for sorting)
    answer_totals = {}
    for answer in unique_answers:
        answer_totals[answer] = int(grouped[grouped['answer'] == answer]['count'].sum())
    
    # Sort answers by total count (descending) for consistent ordering
    sorted_answers = sorted(unique_answers, key=lambda x: answer_totals[x], reverse=True)
    
    # Pivot data
    pivot_data = grouped.pivot(index='answer', columns='year_month', values='percentage').fillna(0)
    months = sorted(pivot_data.columns)
    pivot_data = pivot_data[months]
    
    # Reorder pivot_data rows by sorted_answers
    pivot_data = pivot_data.loc[sorted_answers]
    
    # Color palette for answers
    answer_colors = ['#4285F4', '#EA4335', '#FBBC04', '#34A853', '#FF6D01', '#9334E6']
    
    # Create figure
    fig, ax = plt.subplots(figsize=(16, 10))
    
    # Create stacked bar chart
    bottom = None
    for i, (answer, row) in enumerate(pivot_data.iterrows()):
        ax.bar(months, row.values, bottom=bottom,
              label=answer, color=answer_colors[i % len(answer_colors)], alpha=0.9,
              edgecolor='white', linewidth=1.0)
        if bottom is None:
            bottom = row.values
        else:
            bottom = bottom + row.values
    
    # Formatting
    ax.set_xlabel('Year Month', fontsize=13, fontweight='600', color='#333333')
    ax.set_ylabel('Percentage (%)', fontsize=13, fontweight='600', color='#333333')
    ax.set_title(question_text, fontsize=15, fontweight='600', pad=20, color='#1a1a1a', wrap=True)
    ax.set_ylim(0, 100)
    ax.set_yticks([0, 20, 40, 60, 80, 100])
    ax.set_yticklabels(['0%', '20%', '40%', '60%', '80%', '100%'])
    ax.grid(axis='y', alpha=0.2, linestyle='--', color='#cccccc')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#e0e0e0')
    ax.spines['bottom'].set_color('#e0e0e0')
    
    plt.xticks(rotation=45, ha='right', fontsize=10, color='#666666')
    plt.yticks(fontsize=10, color='#666666')
    
    # Legend
    legend = ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=9,
                       framealpha=0.95, frameon=True, edgecolor='#e0e0e0',
                       facecolor='white')
    legend.get_frame().set_linewidth(0.5)
    
    plt.tight_layout()
    
    # Convert to image bytes
    img_buffer = BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    
    return img_buffer


if __name__ == '__main__':
    import argparse
    
    parser = argparse.ArgumentParser(description='Generate churn trends PowerPoint presentation')
    parser.add_argument('--input', '-i', 
                       default='data/survicate_cancelled_subscriptions_augmented.csv',
                       help='Path to input CSV file')
    parser.add_argument('--output', '-o',
                       default='churn_trends_slides.pptx',
                       help='Path to output PowerPoint file')
    
    args = parser.parse_args()
    
    result = generate_churn_slides(args.input, args.output)
    
    if result:
        sys.exit(0)
    else:
        sys.exit(1)

